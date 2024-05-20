import os
import re
import string
import logging
from pptx import Presentation
import boto3

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_images_from_slide(slide, output_dir, pptx_name):
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == 13:  # Check if shape is an image
            image = shape.image
            image_bytes = image.blob
            image_path = os.path.join(output_dir, f"{pptx_name}_{i+1}.jpg")
            with open(image_path, "wb") as f:
                f.write(image_bytes)
            yield image_path

def extract_text_from_slide(slide):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            cleaned_text = re.sub(r'\W+', ' ', shape.text.strip())
            yield cleaned_text.strip()

def upload_images_to_s3(image_paths, bucket_name, folder_name, aws_access_key_id, aws_secret_access_key):
    s3 = boto3.client('s3', aws_access_key_id=aws_access_key_id, aws_secret_access_key=aws_secret_access_key)
    for path in image_paths:
        filename = os.path.basename(path)
        s3_key = f"{folder_name}/{filename}"
        try:
            with open(path, "rb") as f:
                s3.upload_fileobj(f, bucket_name, s3_key)
        except Exception as e:
            logger.error(f"Failed to upload image {path} to S3: {e}")

def convert_pptx_to_images(pptx_file, bucket_name, folder_name, aws_access_key_id, aws_secret_access_key):
    try:
        # Load PowerPoint presentation
        presentation = Presentation(pptx_file)
        pptx_name = os.path.splitext(os.path.basename(pptx_file))[0]

        # Create output directory if it doesn't exist
        output_dir = os.path.join(os.getcwd(), "temp")
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Initialize lists to store image paths and text content
        image_paths = []
        text_content = []

        # Iterate through each slide in the presentation
        for slide in presentation.slides:
            # Extract text content from slide
            slide_text = extract_text_from_slide(slide)

            # Check if the slide has less than two words or is empty
            if len(list(slide_text)) < 2:
                continue

            # Extract images from slide
            slide_image_paths = extract_images_from_slide(slide, output_dir, pptx_name)
            image_paths.extend(slide_image_paths)

            # Append text content of valid slide
            text_content.extend(slide_text)

        # Upload images to S3 bucket
        upload_images_to_s3(image_paths, bucket_name, folder_name, aws_access_key_id, aws_secret_access_key)

        # Clean up temporary directory
        for path in image_paths:
            os.remove(path)
        os.rmdir(output_dir)

        return text_content

    except Exception as e:
        logger.error(f"An error occurred during pptx to images conversion: {e}")
        return []

# Example usage:
pptx_file = "jpictory.pptx"  # Path to the PowerPoint file
bucket_name = "ppt2video"  # Name of the S3 bucket
folder_name = "proj_pd"  # Name of the folder in the S3 bucket
aws_access_key_id="ASIAUELME625YZYCR4PP"
aws_secret_access_key="66fzCnVjhxSNqmU9yfxDZjcPo0pLNHs8XL65G7Fw"
aws_session_token="IQoJb3JpZ2luX2VjEBUaCXVzLWVhc3QtMiJHMEUCIQCopIFHeyCb1LOA91z2hqyg0YP0H0bXJvNnh/57upa9OQIgXFXuHzNTfjKYG6KFvQ94QVzNeLwcsJ3B1K5MIdaUAukq+QIIbhAAGgwyODQyMzE1MzAxNzEiDEQw6xmB6+pvFNRe7SrWAgnA0uSxZ1DDqpJDTIiI6ISzZHJWVJ+x4oevCID6CZrL/LZdhZgStUIzPhSWH8S6rbzFBtnt1JoZHE9klD0oq3Txjc4wL6VONnXq+hQF1/B55A+Jam0D4fWR1QdyjwgpdpBJAmOKgN8RxUCUFuxzGyiYlg3wWX3A1liiB35Qk4qmMJiARg6AxcFlorsy2s4CWjBEvpjwpuvCv4rYu6kXeBXMvL8RfuZ9vAKM5NV22kiWX5ELrc7+aNckyyEQXTIIwIHxOO4K0r17nknFkQ1bnmX8mRS0612ASE3ZP2U6V4dJn8FWxVlK4nrMpCYrOxFH/kUbwGfLbCG9sTSn3eJwLEbGlPKx88lhZb+KVeDX4eP2wHqdIucxA/YGbKuZ/J31FNp4B+k+C5wZHeNj3Mm6a8yO/livcPzRK7Q3O2FHjwxVWdkxUKoMe+pn0LBjzpV/jPn4IoahXTCF5tGxBjqnAW2ad5yKGlCZdmI8TQ+fEiMzNsE0bmNnkFb7PduEnbBQCjJcdSUlyAYzNUa7IsNzQ25gxgCBhvWSqUf1p03mPXym6xRxyeQQNpui4mtuSECTNXKKX0KVaaT3mDyVaBPycx2+1U/WeUGsokhQopasyXRLUJtp+cWLXjZXZS+JUrWGVIpHtaJIF0POcnGvWl+sZVP5pijFcwCaBY9xqurmucpVW84qGd1D"


text_content = convert_pptx_to_images(pptx_file, bucket_name, folder_name, aws_access_key_id, aws_secret_access_key)

# Print text content of each slide properly with new lines
for i, text in enumerate(text_content):
    logger.info(f"Slide {i+1} Text Content:")
    logger.info(text)
