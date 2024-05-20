import os
import zipfile
import re
from pptx import Presentation
from pptx.util import Inches, Pt

def bk_extract_media_from_pptx(pptx_file, output_dir):
    try:
        # Rename PPTX file to ZIP
        zip_file = os.path.splitext(pptx_file)[0] + ".zip"
        os.rename(pptx_file, zip_file)

        # Open the ZIP file
        with zipfile.ZipFile(zip_file, 'r') as zip_ref:
            # Create output directory if it doesn't exist
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Extract media files from the "ppt/media" folder
            media_paths = {}
            for zip_info in zip_ref.infolist():
                if zip_info.filename.startswith("ppt/media/"):
                    media_path = os.path.join(output_dir, os.path.basename(zip_info.filename))
                    with open(media_path, "wb") as f:
                        f.write(zip_ref.read(zip_info.filename))
                    slide_num = int(re.search(r'\d+', zip_info.filename).group())
                    if slide_num not in media_paths:
                        media_paths[slide_num] = []
                    media_paths[slide_num].append(media_path)

        # Rename the ZIP file back to PPTX
        os.rename(zip_file, pptx_file)
        return media_paths

    except Exception as e:
        print(f"An error occurred during PPTX to media conversion: {e}")
        return {}

def bk_extract_slide_text_and_media(pptx_file, output_dir):
    media_paths = extract_media_from_pptx(pptx_file, output_dir)
    text_content = extract_slide_text(pptx_file)

    # Pair slide text with media paths
    paired_data = {}
    for slide_num, text in text_content.items():
        if slide_num in media_paths:
            paired_data[slide_num] = {"text": text, "media_paths": media_paths[slide_num]}
        else:
            paired_data[slide_num] = {"text": text, "media_paths": []}

    return paired_data

def extract_slide_text(pptx_file):
    try:
        # Load PowerPoint presentation
        presentation = Presentation(pptx_file)
        text_content = {}

        # Extract text content from each slide
        for i, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            text_content[i+1] = "\n".join(slide_text)

        return text_content

    except Exception as e:
        print(f"An error occurred during text extraction: {e}")
        return {}



def extract_media_from_pptx(pptx_file, output_dir):
    try:
        # Rename PPTX file to ZIP
        zip_file = os.path.splitext(pptx_file)[0] + ".zip"
        os.rename(pptx_file, zip_file)

        # Open the ZIP file
        with zipfile.ZipFile(zip_file, 'r') as zip_ref:
            # Create output directory if it doesn't exist
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            # Extract media files from the "ppt/media" folder
            media_paths = {}
            media_positions = {}
            for zip_info in zip_ref.infolist():
                if zip_info.filename.startswith("ppt/media/"):
                    media_path = os.path.join(output_dir, os.path.basename(zip_info.filename))
                    with open(media_path, "wb") as f:
                        f.write(zip_ref.read(zip_info.filename))
                    slide_num = int(re.search(r'\d+', zip_info.filename).group())
                    if slide_num not in media_paths:
                        media_paths[slide_num] = []
                        media_positions[slide_num] = []
                    media_paths[slide_num].append(media_path)

                    # Load the PowerPoint presentation and get the media position
                    presentation = Presentation(pptx_file)
                    slide = presentation.slides[slide_num - 1]
                    for shape in slide.shapes:
                        if hasattr(shape, "image"):
                            if shape.image.blob == zip_ref.read(zip_info.filename):
                                media_positions[slide_num].append({
                                    "left": shape.left,
                                    "top": shape.top,
                                    "width": shape.width,
                                    "height": shape.height
                                })
                                break

        # Rename the ZIP file back to PPTX
        os.rename(zip_file, pptx_file)
        return media_paths, media_positions

    except Exception as e:
        print(f"An error occurred during PPTX to media conversion: {e}")
        return {}, {}

def extract_slide_text_and_media(pptx_file, output_dir):
    media_paths, media_positions = extract_media_from_pptx(pptx_file, output_dir)
    text_content = extract_slide_text(pptx_file)

    # Pair slide text with media paths and positions
    paired_data = {}
    for slide_num, text in text_content.items():
        if slide_num in media_paths:
            paired_data[slide_num] = {
                "text": text,
                "media_info": [{"path": path, "position": position} for path, position in zip(media_paths[slide_num], media_positions[slide_num])]
            }
        else:
            paired_data[slide_num] = {"text": text, "media_info": []}

    return paired_data

# Example usage
pptx_file = "jpictory.pptx"
output_dir = "output_media"
paired_data = extract_slide_text_and_media(pptx_file, output_dir)

# Print paired data (slide number, slide text, media paths)
for slide_num, data in paired_data.items():
    print(f"Slide {slide_num}:")
    print(f"Text:\n{data['text']}")
    print("Media Info:")
    for media_info in data["media_info"]:
        print(f"- Path: {media_info['path']}")
        print(f"  Position: Left={media_info['position']['left']}, Top={media_info['position']['top']}, Width={media_info['position']['width']}, Height={media_info['position']['height']}")
    print()