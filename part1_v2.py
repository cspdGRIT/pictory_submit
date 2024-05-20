import os
import zipfile
import re
from typing import Dict, Tuple, List, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
import logging
import shutil

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def extract_slide_text(pptx_file: str) -> Dict[int, Tuple[List[str], Dict[str, Dict[str, float]]]]:
    """
    Extract the text content and coordinates from each slide in a PowerPoint presentation.

    Args:
        pptx_file (str): The path to the PowerPoint presentation file.

    Returns:
        Dict[int, Tuple[List[str], Dict[str, Dict[str, float]]]]: A dictionary mapping slide numbers to tuples containing a list of text content and a dictionary of text coordinates.
    """
    try:
        presentation = Presentation(pptx_file)
        text_content = {}

        for i, slide in enumerate(presentation.slides):
            slide_text = []
            text_coordinates = {}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                    text_coordinates[shape.text.strip()] = {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
            text_content[i+1] = (slide_text, text_coordinates)

        return text_content
    except Exception as e:
        logger.error(f"An error occurred during text extraction: {e}")
        return {}


def extract_media_from_pptx(pptx_file: str, output_dir: str) -> Tuple[Dict[int, List[str]], Dict[int, List[Dict[str, float]]]]:
    """
    Extract media files (images, videos, etc.) from a PowerPoint presentation and their positions.

    Args:
        pptx_file (str): The path to the PowerPoint presentation file.
        output_dir (str): The directory to save the extracted media files.

    Returns:
        Tuple[Dict[int, List[str]], Dict[int, List[Dict[str, float]]]]:
            - A dictionary mapping slide numbers to lists of media file paths.
            - A dictionary mapping slide numbers to lists of media position information (left, top, width, height).
    """
    try:
        # Create a copy of the original PPTX file
        pptx_copy = os.path.join(output_dir, os.path.basename(pptx_file))
        shutil.copyfile(pptx_file, pptx_copy)

        # Rename the copied PPTX file to ZIP
        zip_file = os.path.splitext(pptx_copy)[0] + ".zip"
        os.rename(pptx_copy, zip_file)

        # Open the ZIP file
        with zipfile.ZipFile(zip_file, 'r') as zip_ref:
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)

            media_paths = {}
            media_positions = {}
            for zip_info in zip_ref.infolist():
                if zip_info.filename.startswith("ppt/media/"):
                    media_path = os.path.join(output_dir, os.path.basename(zip_info.filename))
                    with open(media_path, "wb") as f:
                        f.write(zip_ref.read(zip_info.filename))
                    slide_num = int(re.search(r'\d+', zip_info.filename).group())
                    if slide_num not in media_paths:
                        media_paths[slide_num] = []
                        media_positions[slide_num] = []
                    media_paths[slide_num].append(media_path)

                    presentation = Presentation(pptx_file)
                    try:
                        slide = presentation.slides[slide_num - 1]
                        for shape in slide.shapes:
                            if hasattr(shape, "image"):
                                if shape.image.blob == zip_ref.read(zip_info.filename):
                                    media_positions[slide_num].append({
                                        "left": shape.left,
                                        "top": shape.top,
                                        "width": shape.width,
                                        "height": shape.height
                                    })
                                    break
                    except IndexError:
                        logger.warning(f"Slide {slide_num} not found in the presentation. Skipping media extraction for this slide.")

        # Rename the ZIP file back to PPTX
        os.rename(zip_file, pptx_file)
        return media_paths, media_positions
    except Exception as e:
        logger.error(f"An error occurred during PPTX to media conversion: {e}")
        return {}, {}


def extract_slide_text_and_media(pptx_file: str, output_dir: str) -> Dict[int, Dict[str, Union[Tuple[List[str], Dict[str, Dict[str, float]]], List[Dict[str, float]]]]]:
    """
    Extract the text content, text coordinates, and media files (with their positions) from a PowerPoint presentation.

    Args:
        pptx_file (str): The path to the PowerPoint presentation file.
        output_dir (str): The directory to save the extracted media files.

    Returns:
        Dict[int, Dict[str, Union[Tuple[List[str], Dict[str, Dict[str, float]]], List[Dict[str, float]]]]]:
            A dictionary mapping slide numbers to dictionaries containing the slide text, text coordinates, and media information.
            The text information is a tuple of a list of text content and a dictionary of text coordinates.
            The media information is a list of dictionaries, each with 'path' and 'position' keys.
    """
    media_paths, media_positions = extract_media_from_pptx(pptx_file, output_dir)
    text_content_and_coordinates = extract_slide_text(pptx_file)

    paired_data = {}
    for slide_num, (text, text_coordinates) in text_content_and_coordinates.items():
        media_info = []
        if slide_num in media_paths:
            for path, position in zip(media_paths[slide_num], media_positions[slide_num]):
                media_info.append({"path": path, "position": position})
        paired_data[slide_num] = {
            "text": text,
            "text_coordinates": text_coordinates,
            "media_info": media_info
        }

    return paired_data

# Example usage
pptx_file = "pictory.pptx"  # Path to your PowerPoint presentation file
output_dir = "output_media"  # Directory to save extracted media files
paired_data = extract_slide_text_and_media(pptx_file, output_dir)

for slide_num, data in paired_data.items():
    print(f"Slide {slide_num}:")
    print(f"Text:\n{', '.join(data['text'])}")
    print("Media Info:")
    for media_info in data["media_info"]:
        print(f"- Path: {media_info['path']}")
        print(f"  Position: Left={media_info['position']['left']}, Top={media_info['position']['top']}, Width={media_info['position']['width']}, Height={media_info['position']['height']}")
    print("Text Coordinates:")
    for text, coordinates in data["text_coordinates"].items():
        print(f"- Text: {text}")
        print(f"  Position: Left={coordinates['left']}, Top={coordinates['top']}, Width={coordinates['width']}, Height={coordinates['height']}")
    print()
